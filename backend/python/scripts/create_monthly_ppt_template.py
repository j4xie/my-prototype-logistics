"""Create a 19-slide monthly PPT skeleton for MonthlyPptExporter.

Run once to generate the binary template file. Output:
  smartbi/knowledge/restaurant/ppt_templates/monthly_default.pptx

Each slide has a title text box that the exporter replaces / augments
with data at export time. The exporter treats this as a skeleton — it
can add new text boxes to each slide but doesn't modify the layout.

Matches 鼎鲜 4-1-xx店 - 月度经营分析-24.10.pptx structure (19 slides).
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

# Slide sequence matches 鼎鲜 template structure
SLIDE_TITLES = [
    "月度经营分析",                        # 1. Title page
    "目录",                               # 2. Table of contents
    "一、月度简报",                        # 3. Monthly briefing
    "一、月度经营状况",                    # 4. Section divider
    "1. 1-12月营收完成情况",              # 5. 12-month revenue table
    "2. 5个营业点完成对比",               # 6. 5-venue comparison
    "3. 环比二期对比",                    # 7. 25-row QoQ table
    "4. 厨房档口毛利率",                  # 8. Dept gross margin
    "5. 损溢指标分析",                    # 9. Shrinkage analysis
    "6. 工作改进跟踪表",                  # 10. Work tracking
    "7. 费用开支明细",                    # 11. Expense summary
    "8. 费用科目预算达成",                # 12. Expense subaccount budget
    "9. 人力成本 + 在职人数",             # 13. Labor cost + headcount
    "10. 部门人均产出比",                 # 14. 30-row per-head
    "三、下月计划",                       # 15. Section divider
    "1. 下月营收计划",                    # 16. Next month revenue
    "2. 下月毛利计划",                    # 17. Next month margin
    "3. 下月费用计划",                    # 18. Next month expense
    "4. 下月具体措施",                    # 19. Next month actions
]


def create_template(output_path: Path) -> None:
    """Generate the 19-slide skeleton PPTX."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    for title in SLIDE_TITLES:
        slide_layout = prs.slide_layouts[5]  # blank layout (no placeholders)
        slide = prs.slides.add_slide(slide_layout)

        # Add a title text box at the top
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(12.3)
        height = Inches(0.8)
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.text = title

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


if __name__ == "__main__":
    out = (
        Path(__file__).parents[1]
        / "smartbi"
        / "knowledge"
        / "restaurant"
        / "ppt_templates"
        / "monthly_default.pptx"
    )
    create_template(out)
    print(f"Created: {out}")
    print(f"Size: {out.stat().st_size} bytes")
