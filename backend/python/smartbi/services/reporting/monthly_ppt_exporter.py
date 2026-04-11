"""Monthly PPT exporter — fills the 19-slide template with data.

Reads a skeleton template from knowledge/restaurant/ppt_templates/ and
populates each slide with data from analyzer output (store_pnl_one_pager,
department_pnl, shrinkage_analysis, expense_breakdown, etc.).

Matches 鼎鲜火锅 monthly analysis PPT structure (4-1-xx店 - 月度经营分析
-24.10.pptx, 19 slides). Customer keeps using the same format they've had
for 5+ years; the exporter replaces manual assembly with auto-generation.

Part of P3.5D P2. Section handler + API endpoint come in P3.
"""
from __future__ import annotations

from pathlib import Path
from typing import Any, Optional

from pptx import Presentation
from pptx.util import Inches


class MonthlyPptExporter:
    """Generate a filled monthly PPT from analyzer output.

    Usage:
        exporter = MonthlyPptExporter()
        output_path = exporter.export(
            store_name="鼎鲜火锅·义乌分公司",
            period="2026-02",
            financial_metrics={...},
            diagnostics=[...],
            department_breakdown={...},  # optional
            shrinkage_report={...},       # optional
            expense_breakdown={...},      # optional
            output_path=Path("/tmp/monthly.pptx"),
        )
    """

    def __init__(self, template_path: Optional[Path] = None):
        """Load the template at construction time so export() is fast.

        Defaults to `monthly_default.pptx` in the shipped templates dir.
        Raises FileNotFoundError if template doesn't exist.
        """
        if template_path is None:
            template_path = (
                Path(__file__).parents[2]
                / "knowledge"
                / "restaurant"
                / "ppt_templates"
                / "monthly_default.pptx"
            )
        if not template_path.exists():
            raise FileNotFoundError(
                f"PPT template not found: {template_path}. "
                f"Run scripts/create_monthly_ppt_template.py to generate it."
            )
        self.template_path = template_path

    def export(
        self,
        store_name: str,
        period: str,
        financial_metrics: dict[str, Any],
        diagnostics: list[dict[str, Any]],
        department_breakdown: Optional[dict[str, Any]] = None,
        shrinkage_report: Optional[dict[str, Any]] = None,
        expense_breakdown: Optional[dict[str, Any]] = None,
        output_path: Optional[Path] = None,
    ) -> Path:
        """Generate a filled PPT and save to output_path.

        Optional params (department_breakdown, shrinkage_report,
        expense_breakdown) are populated into their specific slides if
        provided; otherwise those slides remain skeleton-only.

        Returns the path to the saved file.
        """
        prs = Presentation(str(self.template_path))

        # Slide 1: Title page
        self._fill_title_slide(prs.slides[0], store_name, period)

        # Slide 3: Monthly briefing (5 key stats)
        self._fill_briefing_slide(prs.slides[2], financial_metrics)

        # Slide 8: 档口 gross margin (if department_breakdown provided)
        if department_breakdown:
            self._fill_dept_margin_slide(prs.slides[7], department_breakdown)

        # Slide 9: Shrinkage (if shrinkage_report provided)
        if shrinkage_report:
            self._fill_shrinkage_slide(prs.slides[8], shrinkage_report)

        # Slide 12: 25+ expense subaccounts (if expense_breakdown provided)
        if expense_breakdown:
            self._fill_expense_slide(prs.slides[11], expense_breakdown)

        # Slide 14: 30-row per-head productivity (if department_breakdown)
        if department_breakdown:
            self._fill_productivity_slide(prs.slides[13], department_breakdown)

        if output_path is None:
            output_path = Path(f"/tmp/monthly_{store_name}_{period}.pptx")
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        return output_path

    # --- Slide fill helpers ---

    def _fill_title_slide(self, slide, store_name: str, period: str) -> None:
        """Slide 1: replace the generic title with store + period."""
        for shape in slide.shapes:
            if shape.has_text_frame:
                tf = shape.text_frame
                if "月度经营分析" in tf.text:
                    tf.text = f"{store_name}\n月度经营分析\n{period}"
                    return

    def _fill_briefing_slide(self, slide, fm: dict[str, Any]) -> None:
        """Slide 3: 5 key financial stats."""
        text_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(12), Inches(5)
        )
        tf = text_box.text_frame
        tf.text = "本月简报:"
        for label, key in [
            ("营收", "revenue"),
            ("食材成本", "foodCost"),
            ("人力成本", "laborCost"),
            ("毛利率 (折后)", "grossMarginFolded"),
            ("净利润", "netProfit"),
        ]:
            value = fm.get(key)
            if value is None:
                continue
            p = tf.add_paragraph()
            if isinstance(value, (int, float)):
                p.text = f"  {label}: {value:,.2f}"
            else:
                p.text = f"  {label}: {value}"

    def _fill_dept_margin_slide(
        self, slide, dept_breakdown: dict[str, Any]
    ) -> None:
        """Slide 8: department margin breakdown (first 8 departments)."""
        departments = dept_breakdown.get("departments", [])[:8]
        text_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(12), Inches(5)
        )
        tf = text_box.text_frame
        tf.text = "档口毛利率:"
        for dept in departments:
            p = tf.add_paragraph()
            labor = dept.get("laborCost", 0) or 0
            share = dept.get("laborShare", 0) or 0
            p.text = (
                f"  {dept.get('nameZh', dept.get('code', '?'))}: "
                f"人力 {labor:,.0f}, 占比 {share*100:.1f}%"
            )

    def _fill_shrinkage_slide(
        self, slide, shrinkage: dict[str, Any]
    ) -> None:
        """Slide 9: shrinkage per-department + total."""
        total_var = shrinkage.get("totalVarianceAmount", 0) or 0
        text_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(12), Inches(5)
        )
        tf = text_box.text_frame
        tf.text = f"档口损溢总额: {total_var:,.0f}"
        for row in shrinkage.get("rows", [])[:10]:
            p = tf.add_paragraph()
            p.text = (
                f"  {row.get('department', '?')}: "
                f"标准 {row.get('standardCost', 0) or 0:,.0f}, "
                f"实际 {row.get('actualCost', 0) or 0:,.0f}, "
                f"差异 {row.get('varianceAmount', 0) or 0:,.0f}"
            )

    def _fill_expense_slide(
        self, slide, expense_breakdown: dict[str, Any]
    ) -> None:
        """Slide 12: top 15 expense subaccounts."""
        top_accounts = expense_breakdown.get("topAccounts", [])[:15]
        text_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(12), Inches(5)
        )
        tf = text_box.text_frame
        tf.text = "Top 15 费用科目:"
        for acc in top_accounts:
            p = tf.add_paragraph()
            value = acc.get("value", 0) or 0
            p.text = (
                f"  {acc.get('nameZh', acc.get('code', '?'))}: "
                f"{value:,.0f}"
            )

    def _fill_productivity_slide(
        self, slide, dept_breakdown: dict[str, Any]
    ) -> None:
        """Slide 14: 30-row per-head productivity table."""
        departments = dept_breakdown.get("departments", [])[:30]
        text_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(12), Inches(5.5)
        )
        tf = text_box.text_frame
        tf.text = "部门人均产出比:"
        for dept in departments:
            head = dept.get("headCount", 0) or 0
            per_head = dept.get("perHeadCost", 0) or 0
            p = tf.add_paragraph()
            p.text = (
                f"  {dept.get('nameZh', dept.get('code', '?'))}: "
                f"{head} 人, 人均工资 {per_head:,.0f}"
            )
