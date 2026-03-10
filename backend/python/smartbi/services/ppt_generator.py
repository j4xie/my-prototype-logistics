"""Financial PPT Generator — Dynamic slide generation with chart images and AI analysis."""
import base64
import io
import logging
from typing import Dict, Optional, List
from datetime import datetime

logger = logging.getLogger(__name__)

# Chart type display names
CHART_DISPLAY_NAMES = {
    "budget_achievement": "预算完成情况",
    "yoy_mom_comparison": "同比环比分析",
    "pnl_waterfall": "损益表分析",
    "expense_yoy_budget": "费用同比及预算达成",
    "category_yoy_comparison": "各品类同期对比",
    "gross_margin_trend": "毛利率趋势分析",
    "category_structure_donut": "品类结构同比分析",
    "cost_flow_sankey": "成本流向桑基图",
    "variance_analysis": "预算差异分析",
}


class FinancialPPTGenerator:
    """Generate PPTX with dynamic number of chart slides."""

    def generate(self, chart_images: Dict[str, str],
                 analysis_results: Dict[str, str],
                 company_name: str = "白垩纪食品",
                 year: int = 2026,
                 period_type: str = "year",
                 start_month: int = 1, end_month: int = 12,
                 kpi_summary: Optional[Dict] = None,
                 template: str = "default") -> bytes:
        """Generate PPTX bytes from chart images and analysis results."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        except ImportError:
            raise RuntimeError("python-pptx not installed. Run: pip install python-pptx>=0.6.23")

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Colors
        BLUE = RGBColor(0x1B, 0x65, 0xA8)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        DARK = RGBColor(0x2C, 0x3E, 0x50)
        GRAY = RGBColor(0x6B, 0x77, 0x8C)
        GREEN = RGBColor(0x36, 0xB3, 0x7E)

        period_label = self._period_label(year, period_type, start_month, end_month)

        # --- Slide 1: Cover ---
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        # Background
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = BLUE
        # Title
        self._add_text(slide, company_name, Inches(1), Inches(2), Inches(11), Inches(1.2),
                       font_size=Pt(44), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        self._add_text(slide, "财务分析报告", Inches(1), Inches(3.2), Inches(11), Inches(0.8),
                       font_size=Pt(32), color=WHITE, alignment=PP_ALIGN.CENTER)
        self._add_text(slide, period_label, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
                       font_size=Pt(20), color=RGBColor(0xBF, 0xDB, 0xFE), alignment=PP_ALIGN.CENTER)
        self._add_text(slide, f"生成日期: {datetime.now().strftime('%Y-%m-%d')}",
                       Inches(1), Inches(6.5), Inches(11), Inches(0.4),
                       font_size=Pt(14), color=RGBColor(0x93, 0xC5, 0xFD), alignment=PP_ALIGN.CENTER)

        # --- Slide 2: Executive Summary (if KPI summary provided) ---
        if kpi_summary:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._add_text(slide, "执行摘要", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                           font_size=Pt(28), color=BLUE, bold=True)
            # Add KPI cards as text boxes
            kpi_items = kpi_summary.get('items', [])
            for i, kpi in enumerate(kpi_items[:8]):
                col = i % 4
                row = i // 4
                x = Inches(0.5 + col * 3.1)
                y = Inches(1.5 + row * 2.5)
                self._add_kpi_card(slide, kpi, x, y, Inches(2.8), Inches(2.0),
                                   BLUE, WHITE, DARK, GRAY)

        # --- Slides 3..N+2: Chart slides ---
        chart_order = list(CHART_DISPLAY_NAMES.keys())
        # Include charts in defined order, then any extra
        ordered_types = [ct for ct in chart_order if ct in chart_images]
        extra_types = [ct for ct in chart_images if ct not in chart_order]
        ordered_types.extend(extra_types)

        for chart_type in ordered_types:
            img_b64 = chart_images.get(chart_type, '')
            analysis = analysis_results.get(chart_type, '')
            display_name = CHART_DISPLAY_NAMES.get(chart_type, chart_type)

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            # Title bar
            self._add_text(slide, display_name, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                           font_size=Pt(24), color=BLUE, bold=True)

            # Chart image (left 65%)
            if img_b64:
                try:
                    img_data = base64.b64decode(img_b64.split(',')[-1] if ',' in img_b64 else img_b64)
                    img_stream = io.BytesIO(img_data)
                    slide.shapes.add_picture(img_stream, Inches(0.3), Inches(1.0),
                                             Inches(8.5), Inches(5.5))
                except Exception as e:
                    logger.warning(f"Failed to embed chart image for {chart_type}: {e}")
                    self._add_text(slide, f"[图表加载失败: {display_name}]",
                                   Inches(0.3), Inches(3), Inches(8.5), Inches(1),
                                   font_size=Pt(16), color=GRAY)
            else:
                # No image provided — show placeholder
                self._add_text(slide, f"[{display_name} — 图表截图未提供]",
                               Inches(0.3), Inches(3), Inches(8.5), Inches(1),
                               font_size=Pt(16), color=GRAY)

            # Analysis panel (right 35%)
            if analysis:
                self._add_text(slide, "AI 分析", Inches(9.0), Inches(1.0), Inches(4), Inches(0.5),
                               font_size=Pt(16), color=BLUE, bold=True)
                self._add_text(slide, analysis[:800], Inches(9.0), Inches(1.6), Inches(4), Inches(5.0),
                               font_size=Pt(11), color=DARK, word_wrap=True)

        # --- Last Slide: Conclusion ---
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = BLUE
        self._add_text(slide, "谢谢", Inches(1), Inches(2.5), Inches(11), Inches(1.2),
                       font_size=Pt(48), color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        self._add_text(slide, f"{company_name} · {period_label}",
                       Inches(1), Inches(4.0), Inches(11), Inches(0.6),
                       font_size=Pt(20), color=RGBColor(0xBF, 0xDB, 0xFE), alignment=PP_ALIGN.CENTER)
        self._add_text(slide, "由 SmartBI 智能生成",
                       Inches(1), Inches(5.0), Inches(11), Inches(0.4),
                       font_size=Pt(14), color=RGBColor(0x93, 0xC5, 0xFD), alignment=PP_ALIGN.CENTER)

        # Save to bytes
        output = io.BytesIO()
        prs.save(output)
        return output.getvalue()

    def _add_text(self, slide, text, left, top, width, height,
                  font_size=None, color=None, bold=False,
                  alignment=None, word_wrap=True):
        from pptx.util import Pt as PtUtil
        from pptx.enum.text import PP_ALIGN
        if font_size is None:
            font_size = PtUtil(12)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = word_wrap
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = font_size
        if color:
            p.font.color.rgb = color
        p.font.bold = bold
        if alignment:
            p.alignment = alignment
        # Set font to support Chinese
        p.font.name = '微软雅黑'

    def _add_kpi_card(self, slide, kpi, x, y, w, h, blue, white, dark, gray):
        from pptx.util import Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
        # Card background shape
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = white
        shape.line.color.rgb = blue

        label = kpi.get('label', '')
        value = str(kpi.get('value', ''))
        unit = kpi.get('unit', '')

        from pptx.util import Inches as InchesUtil
        self._add_text(slide, label, x, y, w, InchesUtil(0.3),
                       font_size=Pt(12), color=gray, alignment=PP_ALIGN.CENTER)
        self._add_text(slide, f"{value}{unit}", x, y + InchesUtil(0.4), w, InchesUtil(0.4),
                       font_size=Pt(24), color=dark, bold=True, alignment=PP_ALIGN.CENTER)

    def _period_label(self, year, period_type, start_month, end_month):
        if period_type == 'year':
            return f"{year}年度"
        elif period_type == 'month' and start_month == end_month:
            return f"{year}年{start_month}月"
        else:
            return f"{year}年{start_month}-{end_month}月"
