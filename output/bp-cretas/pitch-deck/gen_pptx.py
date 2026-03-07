"""Generate Cretas Pitch Deck PPTX (Detailed — 12 slides)"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SCREENSHOTS_DIR = "C:/Users/Steve/my-prototype-logistics/output/bp-cretas/screenshots"
ASSETS_DIR = "C:/Users/Steve/my-prototype-logistics/output/bp-cretas/pitch-deck/assets"

# Colors — refined palette
DARK = RGBColor(0x11, 0x1d, 0x32)
ACCENT = RGBColor(0xe0, 0x4e, 0x63)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0x64, 0x74, 0x8b)
LIGHT_BG = RGBColor(0xf0, 0xf4, 0xf8)
TABLE_HEADER = RGBColor(0x1e, 0x3a, 0x5f)
ROW_STRIPE = RGBColor(0xf0, 0xf4, 0xf8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    p.font.name = 'Microsoft YaHei'
    return txBox

def add_para(tf, text, size=16, color=DARK, bold=False, align=PP_ALIGN.LEFT, space_before=6):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Microsoft YaHei'
    p.alignment = align
    p.space_before = Pt(space_before)
    return p

def add_metric(slide, left, top, value, label):
    add_text(slide, left, top, 2.5, 0.8, value, size=36, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, left, top+0.7, 2.5, 0.4, label, size=15, color=GRAY, align=PP_ALIGN.CENTER)

def add_table(slide, left, top, width, height, headers, rows):
    tbl_shape = slide.shapes.add_table(len(rows)+1, len(headers), Inches(left), Inches(top), Inches(width), Inches(height))
    tbl = tbl_shape.table
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = WHITE
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.name = 'Microsoft YaHei'
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r+1, c)
            cell.text = str(val)
            # Alternating row stripe
            if r % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ROW_STRIPE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.name = 'Microsoft YaHei'
                p.font.color.rgb = DARK
    return tbl_shape

def add_accent_bar(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.1), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

def add_title_line(slide, left=0.8, top=1.05, width=2.5):
    """Thin accent underline below slide title for visual polish."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

def add_notes(slide, text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text

# ============ SLIDE 1: Cover ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

bg_path = os.path.join(ASSETS_DIR, 'bg_cover.png')
if os.path.exists(bg_path):
    slide.shapes.add_picture(bg_path, Inches(0), Inches(0), Inches(13.333), Inches(7.5))

# Decorative geometric accents (15% opacity)
for cx, cy, sz in [(11.5, 0.8, 1.8), (1.2, 6.2, 1.2), (12.0, 5.5, 0.8)]:
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(sz), Inches(sz))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT
    dot.line.fill.background()
    # Set 15% opacity via XML alpha element
    sp_xml = dot._element
    srgb = sp_xml.find('.//' + qn('a:srgbClr'))
    if srgb is not None:
        alpha_el = srgb.makeelement(qn('a:alpha'), {})
        alpha_el.set('val', '15000')
        srgb.append(alpha_el)

add_text(slide, 2, 1.5, 9, 1.2, "白垩纪", size=60, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 2, 2.7, 9, 0.6, "CRETACEOUS", size=28, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 1.5, 3.5, 10, 0.6, "食品工厂的 AI 管家: 扫码即用, 成本降 90%", size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 1.5, 4.3, 10, 0.5, "让每一家食品企业都拥有自己的 AI 顾问", size=22, color=RGBColor(0xcc,0xcc,0xcc), align=PP_ALIGN.CENTER)
add_text(slide, 4, 5.5, 5, 0.5, "天使轮融资  |  2026", size=22, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, 3, 6.3, 7, 0.4, "AI Agent  |  智能决策  |  自然语言驱动", size=18, color=GRAY, align=PP_ALIGN.CENTER)
add_notes(slide, "开场：白垩纪是食品工厂的AI管家——扫码即用，成本仅为传统ERP的1/10。我们今天寻求天使轮融资。")

# ============ SLIDE 2: Problem ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_text(slide, 0.8, 0.4, 12, 0.7, "中国食品加工业正面临结构性数字化断层", size=32, color=DARK, bold=True)
add_title_line(slide, 0.8, 1.1, 3.0)

add_metric(slide, 1.5, 1.4, "22.7%", "食物损耗浪费率")
add_metric(slide, 4.5, 1.4, "62.6%", "中小企业转型仍处早期")
add_metric(slide, 7.5, 1.4, "3.2%", "实现智能驱动的企业")
add_metric(slide, 10.5, 1.4, "55.6%", "数据孤岛导致决策延迟")

add_text(slide, 0.8, 3.2, 10, 0.5, "现有方案三大结构性失败", size=22, color=ACCENT, bold=True)
add_table(slide, 0.8, 3.9, 7, 1.8,
    ['"用不起"', '"不适用"', '"连不通"'],
    [
        ['平均成本 50-100 万元', '通用软件不懂食品行业', '单一环节数据孤岛'],
        ['实施周期 6-12 个月', '一线不愿用，沦为摆设', '无法业务打通'],
        ['本土化适配差', '仅基础功能，无法提效', '管理统一视图缺失'],
    ])

# Persona story — concrete pain point
add_text(slide, 0.8, 6.0, 7.5, 0.9,
    "典型客户画像: 200人火锅底料厂张厂长，用Excel管30个SKU，每月因保质期漏洞损失~4万。试过金蝶，工人不会用PC端，3个月放弃。",
    size=14, color=TABLE_HEADER, bold=False)
add_text(slide, 0.8, 6.9, 11, 0.4, "七部委 (2025.6): 2027年重点企业80%数字化，2030年全面普及", size=15, color=ACCENT, bold=True)

chart_path = os.path.join(ASSETS_DIR, 'chart_pain_points.png')
if os.path.exists(chart_path):
    slide.shapes.add_picture(chart_path, Inches(7.8), Inches(2.8), Inches(5.3), Inches(3.5))

add_notes(slide, "张厂长的故事就是中国30-40万家食品加工企业的缩影——用Excel管生产、保质期管理靠人记、试过ERP但工人不会用。现有方案要么太贵、要么不适用、要么连不通。七部委政策要求2027年80%数字化，这不是预测，是硬指标。")

# ============ SLIDE 3: Why Now ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_text(slide, 0.8, 0.3, 12, 0.7, "为什么是现在? 五大催化剂正在打开窗口", size=32, color=DARK, bold=True)
add_title_line(slide, 0.8, 1.0, 3.0)

why_now_path = os.path.join(ASSETS_DIR, 'chart_why_now.png')
if os.path.exists(why_now_path):
    slide.shapes.add_picture(why_now_path, Inches(0.8), Inches(1.3), Inches(7), Inches(4))

# Right side: 5 catalyst texts
catalysts_box = slide.shapes.add_textbox(Inches(8.5), Inches(1.3), Inches(4.5), Inches(4.5))
tf = catalysts_box.text_frame
tf.word_wrap = True
add_para(tf, "1. 七部委政策: 2027年80%数字化，30万企业必须上系统", size=16, color=DARK, bold=False, space_before=4)
add_para(tf, "2. AI成本暴降90%+: GPT-4级能力从千元降至分级", size=16, color=DARK, bold=False, space_before=10)
add_para(tf, "3. 食安事件频发: 消费者+监管双向推动溯源需求激增", size=16, color=DARK, bold=False, space_before=10)
add_para(tf, "4. 连锁化率18→25%: 中小企业标准化管理需求爆发", size=16, color=DARK, bold=False, space_before=10)
add_para(tf, "5. 数字化补贴: 地方政府对中小企业上云有专项资金", size=16, color=DARK, bold=False, space_before=10)

add_text(slide, 0.8, 5.8, 12, 0.8, "政策 + 技术 + 市场三重共振 → 2026年是最佳切入窗口", size=18, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_notes(slide, "为什么是现在而不是去年或明年？五大催化剂在2026年同时到位：七部委政策定下2027年80%数字化硬指标，AI API成本暴降让中小企业用得起，食安事件推动溯源刚需，连锁化率提升带来标准化需求，地方补贴降低采购门槛。这是一个确定性的时间窗口。")

# ============ SLIDE 4: Solution ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_text(slide, 0.8, 0.3, 12, 0.7, "白垩纪: 食品企业的 AI Agent，成本仅为传统方案的 1/10", size=32, color=DARK, bold=True)
add_title_line(slide, 0.8, 1.0, 3.0)

add_text(slide, 0.8, 1.1, 6, 0.4, "核心理念: 先收集数据，再成为大脑", size=18, color=TABLE_HEADER, bold=True)
add_table(slide, 0.8, 1.5, 6.2, 2.6,
    ['层级', '定位', '核心价值'],
    [
        ['AI Agent', '智能决策引擎', '自然语言→意图识别(98%+)→自主执行'],
        ['数据采集', 'IoT+移动端', '硬件自动采集+扫码录入，数据真实可信'],
        ['业务管理', 'SaaS 平台', '生产/质量/仓储/财务，移动端零培训'],
        ['餐饮运营', '餐饮垂直模块', '点评雷达图+菜品四象限+门店对比'],
        ['分析洞察', 'SmartBI+LLM', '30秒出分析报告，成本/异常/趋势'],
    ])

arch_path = os.path.join(ASSETS_DIR, 'arch_layers.png')
if os.path.exists(arch_path):
    slide.shapes.add_picture(arch_path, Inches(7.5), Inches(1.1), Inches(5.5), Inches(3.2))

add_text(slide, 0.8, 4.4, 5.5, 0.4, "vs 传统 ERP/MES", size=20, color=ACCENT, bold=True)
add_table(slide, 0.8, 4.8, 6.2, 1.6,
    ['维度', '白垩纪', '传统 ERP/MES'],
    [
        ['成本', 'SaaS 订阅 1-3万/年', '软件许可+实施 50-100万+'],
        ['部署', '10分钟云端上线', '6-12个月实施'],
        ['使用', '移动端扫码，零培训', 'PC端操作，需培训数周'],
        ['分析', 'AI 30秒出成本诊断', '人工汇总1-2天'],
    ])
add_notes(slide, "白垩纪不是又一个ERP——我们是食品工业的AI Agent平台。五大能力层从第一行代码就为食品设计，移动端扫码即用。关键差异：同一AI引擎同时服务工厂和餐饮两个赛道——餐饮模块包含大众点评竞争力雷达图、菜品四象限分析、门店对比。竞品只能做其一，我们覆盖全链。成本仅为传统方案的1/10。")

# ============ SLIDE 5: Market ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_text(slide, 0.8, 0.3, 12, 0.7, "食品数字化是一个 ~2,800 亿元的确定性赛道", size=32, color=DARK, bold=True)
add_title_line(slide, 0.8, 1.0, 3.0)

add_text(slide, 1, 1.3, 5, 0.4, "市场漏斗 (TAM/SAM/SOM)", size=20, color=ACCENT, bold=True)

funnel_path = os.path.join(ASSETS_DIR, 'chart_market_funnel.png')
if os.path.exists(funnel_path):
    slide.shapes.add_picture(funnel_path, Inches(0.5), Inches(1.9), Inches(5.8), Inches(4.5))

add_text(slide, 6.5, 1.3, 6, 0.4, "增长引擎", size=20, color=ACCENT, bold=True)
add_table(slide, 6.5, 1.9, 6, 2.8,
    ['驱动力', '数据'],
    [
        ['制造+AI 市场', '$6B -> $62B (2024-2032, CAGR 35%)'],
        ['AI食品制造', '$95亿→$908亿 (CAGR 28.5%)'],
        ['七部委政策', '2027年重点企业80%数字化'],
        ['目标企业数', '30-40万家食品加工SME (SOM: 首批1万家)'],
        ['餐饮市场(扩展赛道)', '5.57万亿收入, SaaS+AI渗透~30-60亿'],
    ])

add_text(slide, 0.8, 6, 12, 0.5, "海外验证: 核心可比上市公司合计市值 2.21 万亿美元 | 全球农业数字化生态 12.04 万亿美元", size=16, color=TABLE_HEADER, bold=True)
add_notes(slide, "这不是一个需要教育市场的赛道——政策已经替我们做了。七部委联合发文要求2027年80%数字化，30-40万家食品企业必须上系统。我们瞄准的SOM是首批1万家，按效果付费均值约3万/年。餐饮赛道额外提供30-60亿的扩展空间。")

# ============ SLIDE 6: Product Metrics ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_text(slide, 0.8, 0.3, 12, 0.7, "产品已完成 85%+，核心指标全面领先", size=32, color=DARK, bold=True)
add_title_line(slide, 0.8, 1.0, 3.0)

add_metric(slide, 0.8, 1.2, "98%+", "AI意图识别准确率")
add_metric(slide, 3.8, 1.2, "95%", "数据自动采集率")
add_metric(slide, 6.8, 1.2, "80%", "合规留痕覆盖")
add_metric(slide, 9.8, 1.2, "50%", "审计成本降低")

add_text(slide, 0.8, 2.8, 5, 0.4, "技术深度", size=20, color=ACCENT, bold=True)
add_table(slide, 0.8, 3.3, 5.5, 2.2,
    ['指标', '数值'],
    [
        ['API 端点', '1,442 个'],
        ['AI 意图识别', '98%+ (BERT+LLM 混合)'],
        ['全栈覆盖', 'Java+Python+RN+Vue+小程序'],
        ['AI 成本分析', 'Excel->图表->洞察 30秒'],
        ['E2E 测试', '50+ 全部通过'],
    ])

add_text(slide, 7, 2.8, 5.5, 0.4, "产品实拍", size=20, color=ACCENT, bold=True)
screenshots_grid = [
    (7.0, 3.3, "05-web-dashboard.jpeg", "经营驾驶舱"),
    (10.2, 3.3, "07-web-ai-query.jpeg", "AI 智能问答"),
    (7.0, 5.2, "06-web-analysis.jpeg", "智能数据分析"),
    (10.2, 5.2, "10-web-dianping-gap.jpeg", "大众点评竞争力分析"),
]
for left, top, filename, label in screenshots_grid:
    img_path = os.path.join(SCREENSHOTS_DIR, filename)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(3), Inches(1.8))
        add_text(slide, left, top + 1.8, 3, 0.3, label, size=13, color=GRAY, align=PP_ALIGN.CENTER)
add_notes(slide, "我们不是PPT融资——产品已经完成85%以上。1442个API端点、98%的AI识别准确率、全栈5端覆盖。右下角这张截图是大众点评竞争力分析——这是餐饮客户的杀手级功能，帮助餐饮老板一眼看到自己和竞争对手在各维度的差距。")

# ============ SLIDE 6b: Mobile App Demo ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_text(slide, 0.8, 0.3, 12, 0.7, "移动端: 一线工人扫码即用，零培训上手", size=32, color=DARK, bold=True)
add_title_line(slide, 0.8, 1.0, 3.0)

mobile_screenshots = [
    (0.8, 1.3, "01-app-ai-dashboard.png", "AI 工厂仪表盘\n首页AI洞察 + 关键指标"),
    (4.0, 1.3, "02-app-smartbi-kpi.png", "SmartBI 数据分析\n核心KPI + 快捷操作"),
    (7.2, 1.3, "03-app-production.png", "生产批次管理\n实时进度 + 状态追踪"),
    (10.4, 1.3, "04-app-inventory.png", "库存管理\n分类筛选 + 保质期预警"),
]
for left, top, filename, label in mobile_screenshots:
    img_path = os.path.join(SCREENSHOTS_DIR, filename)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(2.6), Inches(5))
        add_text(slide, left, top + 5.1, 2.6, 0.7, label, size=13, color=DARK, align=PP_ALIGN.CENTER)
add_notes(slide, "移动端是我们的核心交互界面。一线工人打开App扫码即可开始使用，零培训成本。AI仪表盘首页直接呈现关键洞察，SmartBI提供数据分析，生产和库存模块覆盖日常操作。")

# ============ SLIDE 8: Traction ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_text(slide, 0.8, 0.3, 12, 0.7, "[X]家食品工厂已在试用，12 个月目标 50 家付费", size=32, color=DARK, bold=True)
add_title_line(slide, 0.8, 1.0, 3.0)

add_text(slide, 0.8, 1.3, 5.5, 0.4, "已完成里程碑", size=20, color=ACCENT, bold=True)
add_table(slide, 0.8, 1.8, 5.5, 3,
    ['时间', '里程碑'],
    [
        ['2025 Q4', '核心后端架构完成, PostgreSQL 迁移'],
        ['2026 Q1', 'AI 意图系统上线 (98% 准确率)'],
        ['2026 Q1', '移动 App 50+ E2E 测试全通过'],
        ['2026 Q2', '[X]家食品工厂试用中 (肉制品/烘焙/调味品)'],
        ['2026 Q2', '客户反馈: "[具体反馈/LOI]"'],
    ])

add_text(slide, 7, 1.3, 5.5, 0.4, "下一阶段目标", size=20, color=ACCENT, bold=True)
add_table(slide, 7, 1.8, 5.5, 2,
    ['时间', '目标', '意义'],
    [
        ['2026 H2', '10-20家付费客户', 'PMF 初步验证'],
        ['2027 H1', 'ARR 100-300万元', '商业模式验证'],
        ['2027 H1', '3个可公开ROI案例', '标杆效应'],
        ['2027 H2', 'Pre-A/A轮融资', '规模化扩张'],
    ])

timeline_path = os.path.join(ASSETS_DIR, 'chart_timeline.png')
if os.path.exists(timeline_path):
    slide.shapes.add_picture(timeline_path, Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.7))

add_text(slide, 0.8, 6.5, 12, 0.4, "GTM 获客路径", size=18, color=ACCENT, bold=True)
add_table(slide, 0.8, 6.9, 11.5, 0.5,
    ['Phase 1: 直销+行业展会', 'Phase 2: 渠道合作+区域代理', 'Phase 3: 生态+口碑裂变'],
    [
        ['目标: 首批50家标杆客户', '行业协会+政府园区合作', '交易平台+品质码消费者触达'],
    ])
add_notes(slide, "目前[X]家食品工厂在试用我们的系统，覆盖肉制品、烘焙、调味品行业。12个月内目标签下50家付费客户，ARR达到100万以上。我们的优势是产品已经ready——不需要再花一年开发，可以直接推向市场验证。")

# ============ SLIDE 9: Business Model ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_text(slide, 0.8, 0.3, 12, 0.7, "Result-as-a-Service: 按效果付费，成本仅为传统方案的 1/10", size=32, color=DARK, bold=True)
add_title_line(slide, 0.8, 1.0, 3.0)

add_text(slide, 0.8, 1.3, 5.5, 0.4, "按效果付费 (20-30% 抽成)", size=20, color=ACCENT, bold=True)
add_table(slide, 0.8, 1.8, 5.5, 2,
    ['收费模式', '比例', '典型场景', '年费参考'],
    [
        ['节省抽成', '20-30%', '降损耗/提效率/减人工', '节省10万→付2-3万'],
        ['增收分成', '20-25%', 'AI分析驱动营收增长', '增收15万→付3-4万'],
        ['保底月费', '500元/月', '覆盖基础运维', '0.6万/年'],
    ])

add_text(slide, 7, 1.3, 5.5, 0.4, "单位经济模型", size=20, color=ACCENT, bold=True)
add_table(slide, 7, 1.8, 5.5, 2.5,
    ['指标', '目标值', '行业基准'],
    [
        ['ARPU (年)', '~3万元', '按效果抽成均值'],
        ['LTV (3年)', '~9万元', '-'],
        ['CAC', '<5,000元', 'SaaS 1-3万'],
        ['LTV/CAC', '~18', '>3 为健康'],
        ['毛利率', '70-80%', 'SaaS 65-85%'],
    ])

# CAC footnote
add_text(slide, 7, 4.4, 5.5, 0.3, "* CAC基于直销+行业展会，天使阶段无付费获客", size=11, color=GRAY)

add_text(slide, 0.8, 4.5, 12, 0.4, "三年营收增长路径", size=20, color=ACCENT, bold=True)

rev_path = os.path.join(ASSETS_DIR, 'chart_revenue_projection.png')
if os.path.exists(rev_path):
    slide.shapes.add_picture(rev_path, Inches(0.8), Inches(5.0), Inches(7), Inches(2.5))

# Financial assumptions (replacing old Y1/Y2/Y3 summary)
fin_box = slide.shapes.add_textbox(Inches(8.2), Inches(5.0), Inches(4.5), Inches(2.5))
tf = fin_box.text_frame
tf.word_wrap = True
add_para(tf, "Base Case 假设推导", size=16, color=ACCENT, bold=True, space_before=2)
add_para(tf, "月新增客户: 3→8→15家", size=14, color=DARK, space_before=6)
add_para(tf, "客户年均节省 ~10万 → 抽成20-30% → ARPU ~3万", size=14, color=DARK, space_before=4)
add_para(tf, "年流失率: ~15% (行业均值20-25%)", size=14, color=DARK, space_before=4)
add_para(tf, "Base Y1: 36家 × 3万 ≈ 110万", size=14, color=TABLE_HEADER, bold=True, space_before=6)
add_para(tf, "Upside Y1: 50家 × 3万 = 150万", size=14, color=DARK, space_before=4)
add_para(tf, "CAC: 直销+行业展会, 天使阶段无付费获客", size=12, color=GRAY, space_before=4)
add_notes(slide, "我们的定价逻辑很简单：帮你省多少/赚多少，我们抽20-30%。典型200人工厂年节省10万以上（损耗降低+效率提升+合规成本减少），我们收取2-3万。客户零风险——不见效不付费，保底月费500元仅覆盖运维。ARPU约3万/年，LTV/CAC约18倍。")

# ============ SLIDE 10: Competition ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_text(slide, 0.8, 0.3, 12, 0.7, '唯一以AI Agent为核心的食品工业智能平台', size=32, color=DARK, bold=True)
add_title_line(slide, 0.8, 1.0, 3.0)

comp_tbl = add_table(slide, 0.8, 1.3, 7.0, 4.0,
    ['能力', '白垩纪', '黑湖智造', '客如云/哗啦啦', '金蝶/用友'],
    [
        ['食品行业专用', '✓', '✗', '✗', '✗'],
        ['中小企业定价', '✓', '部分', '✓', '✗'],
        ['移动端优先', '✓', '部分', '✓', '✗'],
        ['零培训上手', '✓', '✗', '部分', '✗'],
        ['生产管理MES', '✓', '✓', '✗', '✓'],
        ['餐饮运营', '✓', '✗', '✓', '✗'],
        ['AI智能分析', '✓', '✗', '✗', '✗'],
        ['IoT硬件一体', '✓', '✗', '✗', '✗'],
        ['品牌知名度', '低 (初创)', '中', '高', '高'],
        ['线下渠道覆盖', '建设中', '部分', '✓', '✓'],
    ])
tbl = comp_tbl.table
for r in range(tbl.rows.__len__()):
    for c in range(tbl.columns.__len__()):
        for p in tbl.cell(r, c).text_frame.paragraphs:
            p.font.size = Pt(13)

radar_path = os.path.join(ASSETS_DIR, 'chart_competition_radar.png')
if os.path.exists(radar_path):
    slide.shapes.add_picture(radar_path, Inches(8.3), Inches(1.0), Inches(4.5), Inches(4.0))

add_text(slide, 0.8, 5.3, 12, 0.4, "核心护城河", size=20, color=ACCENT, bold=True)
moat = slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(12), Inches(1.5))
tf = moat.text_frame
tf.word_wrap = True
for m in [
    "1. AI Agent架构 — 不是管理工具+AI按钮，而是AI驱动的全新范式",
    "2. 食品行业深度 — 从第一行代码为食品设计，259类意图覆盖全场景",
    "3. 软硬一体+数据飞轮 — 自研IoT确保数据真实，更多工厂→AI更准→更多工厂",
    "4. 双赛道壁垒 — 同一AI Agent服务工厂+餐饮，竞对只做其一"
]:
    add_para(tf, m, size=15, color=DARK, space_before=4)
add_notes(slide, "我们坦诚地列出了弱项——品牌知名度和线下渠道还在建设中。但核心技术能力上，市场上没有任何一家同时覆盖工厂+餐饮、拥有AI原生架构。黑湖智造只做通用制造，客如云只做餐饮收银——都没有AI分析能力。")

# ============ SLIDE 11: Team ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_text(slide, 0.8, 0.3, 12, 0.7, "技术+行业跨学科团队，从一线需求出发推动落地", size=32, color=DARK, bold=True)
add_title_line(slide, 0.8, 1.0, 3.0)

add_text(slide, 0.8, 1.3, 5.5, 0.4, "为什么是我们", size=20, color=ACCENT, bold=True)
why_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.8))
tf = why_box.text_frame
tf.word_wrap = True
add_para(tf, "创始人背景: UCSD AI/ML → Cornell 信息科学，在校期间即开始构建食品行业AI产品", size=15, color=DARK, space_before=2)
add_para(tf, "执行力证明: 一人主导完成 1,442 个 API、5 端全栈、98% AI 意图识别——不是纸上谈兵", size=15, color=DARK, space_before=6)
add_para(tf, "行业连接: 通过家族食品行业资源理解一线需求，产品设计源于真实工厂调研", size=15, color=DARK, space_before=6)
add_para(tf, "", size=6, space_before=0)
add_para(tf, "坦诚的 Gap:", size=15, color=ACCENT, bold=True, space_before=4)
add_para(tf, "缺 B2B SaaS 销售经验 → 天使轮首要补充销售负责人", size=14, color=GRAY, space_before=2)
add_para(tf, "缺资深食品行业 PM → 招聘计划中优先级 #1", size=14, color=GRAY, space_before=2)

add_text(slide, 7, 1.3, 5.5, 0.4, "技术实力证明", size=20, color=ACCENT, bold=True)
add_metric(slide, 7.2, 1.9, "1,442", "API 端点")
add_metric(slide, 9.7, 1.9, "98%+", "AI 识别准确率")

add_text(slide, 0.8, 4.5, 5.5, 0.4, "天使轮后招聘计划", size=20, color=ACCENT, bold=True)
add_table(slide, 0.8, 5.0, 5.5, 1.8,
    ['角色', '人数', '目的'],
    [
        ['食品行业产品经理', '1-2人', '行业know-how补充'],
        ['销售负责人', '1人', 'B2B SaaS获客体系'],
        ['客户成功经理', '1-2人', '标杆客户深度服务'],
    ])

add_text(slide, 7, 4.5, 5.5, 0.4, "核心团队", size=20, color=ACCENT, bold=True)

# Photo placeholders (gray rectangles)
for i, x in enumerate([7.0, 8.5]):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(5.0), Inches(1.2), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xdd, 0xdd, 0xdd)
    shape.line.color.rgb = RGBColor(0xbb, 0xbb, 0xbb)
    # Add [照片] label inside
    add_text(slide, x + 0.15, 5.5, 0.9, 0.4, "[照片]", size=12, color=GRAY, align=PP_ALIGN.CENTER)

team_box = slide.shapes.add_textbox(Inches(10.0), Inches(5.0), Inches(3.2), Inches(1.8))
tf = team_box.text_frame
tf.word_wrap = True
add_para(tf, "谢杰涛 - 创始人 & CEO", size=16, color=DARK, bold=True, space_before=2)
add_para(tf, "UCSD AI/ML + Cornell AI/信息科学", size=14, color=DARK, space_before=2)
add_para(tf, "独立主导全栈产品 (1,442 API)", size=14, color=GRAY, space_before=2)
add_para(tf, "", size=4, color=GRAY, space_before=0)
add_para(tf, "核心工程师: Cornell / UC Berkeley CS", size=14, color=DARK, space_before=2)

# Advisor section
add_text(slide, 7, 6.8, 2, 0.3, "顾问 Advisory", size=18, color=ACCENT, bold=True)
add_text(slide, 9.2, 6.8, 4, 0.3, "[顾问姓名] — [背景/行业资源]", size=14, color=DARK)
add_notes(slide, "我们是一支技术能力极强的团队，独立完成了全栈产品开发。天使轮后核心补充食品行业经验者和销售能力。顾问席位预留给带行业资源的产业投资人或行业专家。")

# ============ SLIDE 12: The Ask ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text(slide, 1, 0.5, 11, 0.6, "天使轮融资: 800-1,200 万元", size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, 0.5, 1.5, 4, 0.8, "800-1,200万", size=36, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 0.5, 2.3, 4, 0.4, "融资金额 (RMB)", size=15, color=RGBColor(0xaa,0xaa,0xaa), align=PP_ALIGN.CENTER)

add_text(slide, 4.8, 1.5, 3.5, 0.8, "12-18%", size=36, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 4.8, 2.3, 3.5, 0.4, "出让股权", size=15, color=RGBColor(0xaa,0xaa,0xaa), align=PP_ALIGN.CENTER)

add_text(slide, 8.8, 1.5, 4.2, 0.8, "5,000-8,000万", size=36, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 8.8, 2.3, 4.2, 0.4, "投后估值 (RMB)", size=15, color=RGBColor(0xaa,0xaa,0xaa), align=PP_ALIGN.CENTER)

add_table(slide, 1, 3.2, 5.5, 2,
    ['用途', '比例', '金额'],
    [
        ['产品研发', '40%', '400万'],
        ['销售市场', '30%', '300万'],
        ['团队扩充', '20%', '200万'],
        ['运营储备', '10%', '100万'],
    ])

add_table(slide, 7, 3.2, 5.5, 2,
    ['时间', '里程碑', '意义'],
    [
        ['3个月', '5家测试转付费', 'PMF信号'],
        ['6个月', '20家+2个ROI案例', '市场验证'],
        ['12个月', '50家+ARR 100万', 'Pre-A/A轮'],
    ])

fund_path = os.path.join(ASSETS_DIR, 'chart_fund_allocation.png')
if os.path.exists(fund_path):
    slide.shapes.add_picture(fund_path, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.2))

add_text(slide, 1, 6.2, 11, 0.4, "15分钟 Demo 即可见真章 | WeChat: [微信号] | [email]", size=18, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 1, 6.6, 11, 0.4, "下一轮: Pre-A / A 轮  3,000-5,000 万元 (12个月后)", size=16, color=RGBColor(0xaa,0xaa,0xaa), align=PP_ALIGN.CENTER)
add_text(slide, 2, 7.0, 9, 0.4, "白垩纪  |  食品工厂的 AI 管家  |  让每一家食品企业都拥有自己的AI顾问", size=14, color=RGBColor(0x99,0x99,0x99), align=PP_ALIGN.CENTER)
add_notes(slide, "我们寻求800-1200万天使轮融资，出让12-18%股权。核心目标：12个月内签约50家付费客户，ARR达到100万以上，产出3个可公开的ROI案例。欢迎15分钟Demo，让产品说话。我们优先寻求产业型投资人——带食品行业资源的投资人，比纯财务投资人更有价值。")

# Save
output_path = "C:/Users/Steve/my-prototype-logistics/output/bp-cretas/pitch-deck/pitch-deck.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
