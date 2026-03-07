"""Generate Cretas Pitch Deck LITE (10 slides, mobile-optimized, cold-send)

Design principles:
- Max 40 words per slide
- Min font size 20pt
- No tables — large metrics + charts only
- Assertion-Evidence framework: 1 headline + 1 visual per slide
- Guy Kawasaki 10/20/30 rule
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SCREENSHOTS_DIR = "C:/Users/Steve/my-prototype-logistics/output/bp-cretas/screenshots"
ASSETS_DIR = "C:/Users/Steve/my-prototype-logistics/output/bp-cretas/pitch-deck/assets"

# Colors — refined palette
DARK = RGBColor(0x11, 0x1d, 0x32)
ACCENT = RGBColor(0xe0, 0x4e, 0x63)
WHITE = RGBColor(0xff, 0xff, 0xff)
GRAY = RGBColor(0x64, 0x74, 0x8b)
LIGHT_GRAY = RGBColor(0x94, 0xa3, 0xb8)
TABLE_HEADER = RGBColor(0x1e, 0x3a, 0x5f)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=20, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    p.font.name = 'Microsoft YaHei'
    return txBox


def add_big_metric(slide, left, top, value, label, value_size=48, label_size=20):
    """Large metric for lite deck — bigger than detailed version."""
    add_text(slide, left, top, 3.5, 1.0, value, size=value_size, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, left, top + 0.9, 3.5, 0.5, label, size=label_size, color=GRAY, align=PP_ALIGN.CENTER)


def add_accent_bar(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.1), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

def add_title_line(slide, left=0.8, top=1.05, width=2.5):
    """Thin accent underline below slide title."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_img(slide, filename, left, top, width, height, from_assets=True):
    """Add image if it exists. Returns True if added."""
    base = ASSETS_DIR if from_assets else SCREENSHOTS_DIR
    path = os.path.join(base, filename)
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
        return True
    print(f'  WARNING: missing {path}')
    return False


def add_notes(slide, text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ============ SLIDE 1: Cover ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_img(slide, 'bg_cover.png', 0, 0, 13.333, 7.5)

# Decorative geometric accents (15% opacity)
for cx, cy, sz in [(11.5, 0.8, 1.8), (1.2, 6.2, 1.2), (12.0, 5.5, 0.8)]:
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(sz), Inches(sz))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT
    dot.line.fill.background()
    sp_xml = dot._element
    srgb = sp_xml.find('.//' + qn('a:srgbClr'))
    if srgb is not None:
        alpha_el = srgb.makeelement(qn('a:alpha'), {})
        alpha_el.set('val', '15000')
        srgb.append(alpha_el)

add_text(slide, 2, 1.5, 9, 1.2, "白垩纪", size=60, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 2, 2.7, 9, 0.6, "CRETACEOUS", size=28, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 1.5, 3.5, 10, 0.6, "食品工厂的 AI 管家: 扫码即用, 成本降 90%", size=28, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 1.5, 4.3, 10, 0.5, "让每一家食品企业都拥有自己的 AI 顾问", size=22, color=RGBColor(0xcc, 0xcc, 0xcc), align=PP_ALIGN.CENTER)
add_text(slide, 4, 5.5, 5, 0.5, "天使轮融资  |  2026", size=22, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, 3, 6.3, 7, 0.4, "AI Agent  |  智能决策  |  自然语言驱动", size=20, color=GRAY, align=PP_ALIGN.CENTER)
add_notes(slide, "白垩纪——食品工厂的AI管家，扫码即用，成本降90%。天使轮融资800-1200万。")

# ============ SLIDE 2: Problem ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_text(slide, 0.8, 0.4, 12, 0.7, "30-40万家食品企业必须数字化 — 七部委政策驱动", size=32, color=DARK, bold=True)
add_title_line(slide, 0.8, 1.1, 3.0)

add_big_metric(slide, 0.5, 1.8, "22.7%", "食物损耗浪费率")
add_big_metric(slide, 4.5, 1.8, "62.6%", "仍处数字化早期")
add_big_metric(slide, 8.5, 1.8, "3.2%", "实现智能驱动")

add_img(slide, 'chart_pain_points.png', 7.5, 3.8, 5.5, 3.5)
add_text(slide, 0.8, 3.8, 6.5, 0.8, "现有方案: 用不起 · 不适用 · 连不通", size=26, color=TABLE_HEADER, bold=True)
add_text(slide, 0.8, 4.8, 6.5, 0.8, "200人底料厂张厂长:\nExcel管30个SKU, 月损4万, 试过金蝶放弃了", size=20, color=DARK)
add_text(slide, 0.8, 6.0, 6, 0.6, "七部委: 2027年 80% 数字化", size=22, color=ACCENT, bold=True)
add_notes(slide, "张厂长的故事是30-40万家食品企业的缩影。现有方案太贵不适用，七部委政策已定硬指标。")

# ============ SLIDE 3: Why Now ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_text(slide, 0.8, 0.3, 12, 0.7, "五大催化剂打开 2026 窗口", size=36, color=DARK, bold=True)
add_title_line(slide, 0.8, 1.05, 3.0)

add_img(slide, 'chart_why_now.png', 0.5, 1.3, 7.5, 4.5)

# Right side: 3 large catalyst highlights
add_text(slide, 8.5, 1.5, 4.5, 1.0, "七部委政策\n2027年80%数字化", size=26, color=DARK, bold=True)
add_text(slide, 8.5, 3.0, 4.5, 1.0, "AI成本暴降90%+\nGPT-4级能力平民化", size=26, color=DARK, bold=True)
add_text(slide, 8.5, 4.5, 4.5, 1.0, "连锁化率18→25%\n标准化需求激增", size=26, color=DARK, bold=True)

add_text(slide, 0.8, 6.3, 12, 0.5, "政策 + 技术 + 市场三重共振 = 2026 最佳切入窗口", size=24, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_notes(slide, "五大催化剂在2026年同时到位：政策硬指标、AI成本暴降、食安刚需、连锁化、补贴。")

# ============ SLIDE 4: Solution ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_text(slide, 0.8, 0.3, 12, 0.7, "AI Agent, 不是又一个 ERP — 成本 1/10", size=32, color=DARK, bold=True)
add_title_line(slide, 0.8, 1.0, 3.0)

add_img(slide, 'arch_layers.png', 1.0, 1.5, 6.5, 4.5)

# Right side: 3 contrast metrics
add_text(slide, 8.0, 1.5, 5, 0.6, "1-3万/年", size=40, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 8.0, 2.2, 5, 0.4, "vs 传统 50-100万", size=20, color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, 8.0, 3.2, 5, 0.6, "10 分钟上线", size=40, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 8.0, 3.9, 5, 0.4, "vs 6-12 个月实施", size=20, color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, 8.0, 4.9, 5, 0.6, "30秒 AI 分析", size=40, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 8.0, 5.6, 5, 0.4, "vs 人工 1-2 天", size=20, color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, 0.8, 6.5, 12, 0.5, "自然语言 → 意图识别(98%+) → 自主执行 → 结果", size=22, color=TABLE_HEADER, bold=True, align=PP_ALIGN.CENTER)
add_notes(slide, "白垩纪是AI Agent，不是ERP。成本1/10，部署10分钟，AI 30秒出分析。")

# ============ SLIDE 5: Market ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_text(slide, 0.8, 0.3, 12, 0.7, "~2,800 亿确定性赛道", size=36, color=DARK, bold=True)
add_title_line(slide, 0.8, 1.05, 3.0)

add_img(slide, 'chart_market_funnel.png', 0.5, 1.3, 6.5, 5.0)

add_big_metric(slide, 7.5, 1.5, "~2,800亿", "TAM 工业软件", value_size=42)
add_big_metric(slide, 7.5, 3.0, "~100亿", "SAM 食品MES+AI", value_size=42)
add_big_metric(slide, 7.5, 4.5, "3.0亿", "SOM 首批1万家", value_size=42)

add_text(slide, 0.8, 6.6, 12, 0.5, "七部委政策: 2027年重点企业 80% 数字化 | CAGR 28-35%", size=22, color=TABLE_HEADER, bold=True, align=PP_ALIGN.CENTER)
add_notes(slide, "TAM 2800亿工业软件，SAM 100亿食品MES+AI，SOM首批1万家×3万=3亿。")

# ============ SLIDE 6: Product + Mobile ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_text(slide, 0.8, 0.3, 12, 0.7, "85%+ 完成，全端覆盖", size=32, color=DARK, bold=True)
add_title_line(slide, 0.8, 1.0, 3.0)

add_big_metric(slide, 0.3, 1.3, "98%+", "AI 意图识别", value_size=44)
add_big_metric(slide, 3.5, 1.3, "95%", "数据自动采集", value_size=44)
add_big_metric(slide, 6.7, 1.3, "1,442", "API 端点", value_size=44)
add_big_metric(slide, 9.9, 1.3, "50+", "E2E 测试通过", value_size=44)

# Web screenshots: 3 in a row (larger, cleaner)
web_screenshots = [
    (0.5, 3.0, "05-web-dashboard.jpeg", "经营驾驶舱"),
    (4.7, 3.0, "07-web-ai-query.jpeg", "AI 智能问答"),
    (8.9, 3.0, "10-web-dianping-gap.jpeg", "点评竞争力分析"),
]
for left, top, filename, label in web_screenshots:
    add_img(slide, filename, left, top, 3.5, 2.2, from_assets=False)
    add_text(slide, left, top + 2.25, 3.5, 0.4, label, size=20, color=DARK, align=PP_ALIGN.CENTER)

# Single line instead of mobile thumbnails (Solution page already covers mobile)
add_text(slide, 0.8, 6.0, 12, 0.5, "全端覆盖: Web + App + 微信小程序 | 移动端扫码即用, 零培训", size=22, color=TABLE_HEADER, bold=True, align=PP_ALIGN.CENTER)
add_notes(slide, "产品85%+完成，Web+App+小程序全端覆盖。截图均为实际运行产品。")

# ============ SLIDE 7: Traction ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_text(slide, 0.8, 0.3, 12, 0.7, "[X]家工厂试用中，12 个月目标 50 家", size=32, color=DARK, bold=True)
add_title_line(slide, 0.8, 1.0, 3.0)

add_img(slide, 'chart_timeline.png', 0.5, 1.3, 12.3, 2.8)

# Customer signal highlight
add_text(slide, 0.8, 4.2, 12, 0.6, "[X]家试用 (肉制品/烘焙/调味品) | [试用反馈关键词]", size=24, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

add_big_metric(slide, 0.8, 5.1, "2026 H2", "10-20家付费客户", value_size=36)
add_big_metric(slide, 4.8, 5.1, "2027 H1", "ARR 100万+", value_size=36)
add_big_metric(slide, 8.8, 5.1, "2027 H2", "Pre-A / A轮", value_size=36)

add_text(slide, 0.8, 6.6, 12, 0.5, "GTM: 直销+行业展会 → 渠道合作 → 生态裂变", size=22, color=TABLE_HEADER, bold=True, align=PP_ALIGN.CENTER)
add_notes(slide, "[X]家食品工厂试用中，12个月目标50家付费。GTM：直销→渠道→裂变。")

# ============ SLIDE 8: Business Model ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_text(slide, 0.8, 0.3, 12, 0.7, "按效果付费: 节省/增收的 20-30% 抽成", size=32, color=DARK, bold=True)
add_title_line(slide, 0.8, 1.0, 3.0)

add_img(slide, 'chart_revenue_projection.png', 0.5, 1.3, 7.0, 3.5)

add_text(slide, 8.0, 1.5, 5, 0.7, "LTV/CAC", size=28, color=DARK, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 8.0, 2.1, 5, 0.7, "~18", size=44, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 8.0, 2.8, 5, 0.4, "远超行业3倍健康线", size=18, color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, 8.0, 3.5, 5, 0.7, "ARPU ~3万", size=36, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 8.0, 4.1, 5, 0.4, "节省10万 × 抽成25%", size=18, color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, 8.0, 4.7, 5, 0.7, "毛利 70-80%", size=36, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 8.0, 5.3, 5, 0.4, "客户零风险，不见效不付费", size=18, color=GRAY, align=PP_ALIGN.CENTER)

# Assumptions derivation (replaces old pricing detail)
add_text(slide, 0.8, 5.5, 12, 0.5, "Base: 月增3家 × 3万 ≈ Y1 110万 | Upside: Y1 150万 | Y3 盈利", size=24, color=TABLE_HEADER, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 0.8, 6.2, 12, 0.5, "* CAC: 直销+展会, 天使阶段无付费获客", size=20, color=GRAY, align=PP_ALIGN.CENTER)
add_notes(slide, "按效果付费：帮客户省/赚多少抽20-30%。客户年均节省10万→ARPU约3万。月增3家→Y1约110万。客户零风险，不见效不付费。")

# ============ SLIDE 9: Competition ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_text(slide, 0.8, 0.3, 12, 0.7, "唯一 AI Agent + 食品专用", size=36, color=DARK, bold=True)
add_title_line(slide, 0.8, 1.05, 3.0)

add_img(slide, 'chart_competition_quadrant.png', 1.5, 1.2, 7.5, 5.0)

# Right side: moat summary
add_text(slide, 9.5, 1.5, 3.5, 0.5, "护城河", size=24, color=ACCENT, bold=True)
add_text(slide, 9.5, 2.2, 3.5, 0.8, "AI Agent 架构\n非管理工具+AI按钮", size=20, color=DARK, bold=True)
add_text(slide, 9.5, 3.4, 3.5, 0.8, "食品深度\n259类意图全覆盖", size=20, color=DARK, bold=True)
add_text(slide, 9.5, 4.6, 3.5, 0.8, "双赛道壁垒\n工厂+餐饮同一引擎", size=20, color=DARK, bold=True)
add_text(slide, 9.5, 5.8, 3.5, 0.8, "数据飞轮\n更多工厂→AI更准", size=20, color=DARK, bold=True)
add_notes(slide, "唯一AI Agent+食品专用平台。四大护城河：AI架构、食品深度、双赛道、数据飞轮。")

# ============ SLIDE 10: The Ask + Team + CTA ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_text(slide, 1, 0.5, 11, 0.8, "天使轮融资", size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, 0.3, 1.8, 4.2, 0.9, "800-1,200万", size=40, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 0.3, 2.7, 4.2, 0.5, "融资金额 (RMB)", size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(slide, 4.6, 1.8, 4, 0.9, "12-18%", size=40, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 4.6, 2.7, 4, 0.5, "出让股权", size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(slide, 8.8, 1.8, 4.3, 0.9, "5,000-8,000万", size=36, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, 8.8, 2.7, 4.3, 0.5, "投后估值", size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
# Fund usage one-liner
add_text(slide, 0.8, 3.5, 12, 0.5, "资金用途: 产品研发 40% | 销售市场 30% | 团队 20% | 储备 10%", size=22, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(slide, 0.8, 4.3, 12, 0.6, "12个月: 50家客户 | ARR 100万+ | Pre-A/A轮", size=28, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# Team
add_text(slide, 0.8, 5.2, 12, 0.5, "谢杰涛 | UCSD + Cornell AI/ML | 独立主导全栈产品 (1,442 API)", size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# CTA
add_text(slide, 0.8, 5.9, 12, 0.5, "15分钟 Demo 即可见真章 | WeChat: [微信号] | [email]", size=22, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, 2, 6.6, 9, 0.4, "白垩纪  |  食品工厂的 AI 管家  |  让每一家食品企业都拥有自己的AI顾问", size=20, color=RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.CENTER)
add_notes(slide, "天使轮800-1200万，12-18%股权。12个月50家客户+ARR100万。15分钟Demo见真章。")

# Save
output_path = "C:/Users/Steve/my-prototype-logistics/output/bp-cretas/pitch-deck/pitch-deck-lite.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
